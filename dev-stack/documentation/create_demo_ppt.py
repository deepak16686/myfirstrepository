"""
AI DevOps Pipeline Platform - Demo Presentation Generator
Creates a professional PowerPoint with animations for demoing:
  - GitLab Pipeline Generator
  - Jenkins Pipeline Generator
  - GitHub Actions Pipeline Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
import copy
from lxml import etree
import os

# ── Color Palette (Light Theme — Orange Accent) ───────────────
COLORS = {
    "bg_dark":       RGBColor(0xFF, 0xFF, 0xFF),   # White background
    "bg_card":       RGBColor(0xFB, 0xF5, 0xEF),   # Warm cream card
    "accent_blue":   RGBColor(0xE8, 0x6A, 0x17),   # Orange primary
    "accent_green":  RGBColor(0x2E, 0x7D, 0x32),   # Forest green
    "accent_orange": RGBColor(0xF5, 0x7C, 0x00),   # Deep orange
    "accent_red":    RGBColor(0xC6, 0x28, 0x28),   # Deep red
    "accent_purple": RGBColor(0x6A, 0x1B, 0x9A),   # Deep purple
    "text_white":    RGBColor(0x1A, 0x1A, 0x2E),   # Near-black text (on white bg)
    "text_light":    RGBColor(0x4A, 0x4A, 0x5A),   # Secondary text
    "text_dim":      RGBColor(0x8A, 0x8A, 0x9A),   # Muted text
    "gitlab_orange": RGBColor(0xE8, 0x5D, 0x04),   # GitLab brand orange
    "jenkins_red":   RGBColor(0xCC, 0x33, 0x33),   # Jenkins brand red
    "github_blue":   RGBColor(0x0D, 0x69, 0xD0),   # GitHub Actions brand blue
    "border":        RGBColor(0xE0, 0xD5, 0xC8),   # Warm border
    "white":         RGBColor(0xFF, 0xFF, 0xFF),   # Actual white (for text on colored fills)
    "orange_light":  RGBColor(0xFF, 0xF3, 0xE0),   # Light orange tint
    "title_bar":     RGBColor(0xE8, 0x6A, 0x17),   # Orange title bar
}

# ── Animation Helpers ──────────────────────────────────────────

def _add_anim_namespace(slide):
    """Ensure slide has timing/animation elements."""
    sld = slide._element
    timing = sld.find(qn('p:timing'))
    if timing is None:
        timing = etree.SubElement(sld, qn('p:timing'))
    tn_lst = timing.find(qn('p:tnLst'))
    if tn_lst is None:
        tn_lst = etree.SubElement(timing, qn('p:tnLst'))
    par = tn_lst.find(qn('p:par'))
    if par is None:
        par = etree.SubElement(tn_lst, qn('p:par'))
        c_tn = etree.SubElement(par, qn('p:cTn'))
        c_tn.set('id', '1')
        c_tn.set('dur', 'indefinite')
        c_tn.set('restart', 'never')
        c_tn.set('nodeType', 'tmRoot')
        child_tn_lst = etree.SubElement(c_tn, qn('p:childTnLst'))
    else:
        c_tn = par.find(qn('p:cTn'))
        child_tn_lst = c_tn.find(qn('p:childTnLst'))
        if child_tn_lst is None:
            child_tn_lst = etree.SubElement(c_tn, qn('p:childTnLst'))
    return child_tn_lst, timing


def _get_next_id(timing):
    """Get next unique animation ID."""
    max_id = 1
    for el in timing.iter():
        ctn_id = el.get('id')
        if ctn_id and ctn_id.isdigit():
            max_id = max(max_id, int(ctn_id))
    return max_id + 1


def add_fade_in(slide, shape, delay_ms=0, duration_ms=500, trigger="afterPrev"):
    """Add fade-in entrance animation to a shape."""
    child_tn_lst, timing = _add_anim_namespace(slide)
    nid = _get_next_id(timing)
    sp_id = shape.shape_id

    seq = child_tn_lst.find(qn('p:seq'))
    if seq is None:
        seq = etree.SubElement(child_tn_lst, qn('p:seq'))
        seq.set('concurrent', '1')
        seq.set('nextAc', 'seek')
        seq_ctn = etree.SubElement(seq, qn('p:cTn'))
        seq_ctn.set('id', str(nid)); nid += 1
        seq_ctn.set('dur', 'indefinite')
        seq_ctn.set('nodeType', 'mainSeq')
        seq_child = etree.SubElement(seq_ctn, qn('p:childTnLst'))
        # prev/next conditions
        prev_cond = etree.SubElement(seq, qn('p:prevCondLst'))
        pc = etree.SubElement(prev_cond, qn('p:cond'))
        pc.set('evt', 'onPrev')
        pc.set('delay', '0')
        pc_tgt = etree.SubElement(pc, qn('p:tgtEl'))
        etree.SubElement(pc_tgt, qn('p:sldTgt'))
        next_cond = etree.SubElement(seq, qn('p:nextCondLst'))
        nc = etree.SubElement(next_cond, qn('p:cond'))
        nc.set('evt', 'onNext')
        nc.set('delay', '0')
        nc_tgt = etree.SubElement(nc, qn('p:tgtEl'))
        etree.SubElement(nc_tgt, qn('p:sldTgt'))
    else:
        seq_ctn = seq.find(qn('p:cTn'))
        seq_child = seq_ctn.find(qn('p:childTnLst'))

    # Create animation group
    par1 = etree.SubElement(seq_child, qn('p:par'))
    ctn1 = etree.SubElement(par1, qn('p:cTn'))
    ctn1.set('id', str(nid)); nid += 1
    ctn1.set('fill', 'hold')

    st_cond = etree.SubElement(ctn1, qn('p:stCondLst'))
    cond = etree.SubElement(st_cond, qn('p:cond'))
    cond.set('delay', '0')

    child1 = etree.SubElement(ctn1, qn('p:childTnLst'))
    par2 = etree.SubElement(child1, qn('p:par'))
    ctn2 = etree.SubElement(par2, qn('p:cTn'))
    ctn2.set('id', str(nid)); nid += 1
    ctn2.set('fill', 'hold')

    st_cond2 = etree.SubElement(ctn2, qn('p:stCondLst'))
    cond2 = etree.SubElement(st_cond2, qn('p:cond'))
    cond2.set('delay', str(delay_ms))

    child2 = etree.SubElement(ctn2, qn('p:childTnLst'))
    par3 = etree.SubElement(child2, qn('p:par'))
    ctn3 = etree.SubElement(par3, qn('p:cTn'))
    ctn3.set('id', str(nid)); nid += 1
    ctn3.set('presetID', '10')      # fade
    ctn3.set('presetClass', 'entr')
    ctn3.set('presetSubtype', '0')
    ctn3.set('fill', 'hold')
    ctn3.set('nodeType', 'afterEffect')

    st_cond3 = etree.SubElement(ctn3, qn('p:stCondLst'))
    cond3 = etree.SubElement(st_cond3, qn('p:cond'))
    cond3.set('delay', '0')

    child3 = etree.SubElement(ctn3, qn('p:childTnLst'))

    # Fade effect (animEffect)
    anim_effect = etree.SubElement(child3, qn('p:animEffect'))
    anim_effect.set('transition', 'in')
    anim_effect.set('filter', 'fade')
    ae_ctn = etree.SubElement(anim_effect, qn('p:cBhvr'))
    ae_ctn_inner = etree.SubElement(ae_ctn, qn('p:cTn'))
    ae_ctn_inner.set('id', str(nid)); nid += 1
    ae_ctn_inner.set('dur', str(duration_ms))
    tgt_el = etree.SubElement(ae_ctn, qn('p:tgtEl'))
    sp_tgt = etree.SubElement(tgt_el, qn('p:spTgt'))
    sp_tgt.set('spid', str(sp_id))

    # Set element (make visible)
    set_el = etree.SubElement(child3, qn('p:set'))
    set_ctn = etree.SubElement(set_el, qn('p:cBhvr'))
    set_ctn_inner = etree.SubElement(set_ctn, qn('p:cTn'))
    set_ctn_inner.set('id', str(nid)); nid += 1
    set_ctn_inner.set('dur', '1')
    set_ctn_inner.set('fill', 'hold')

    set_st = etree.SubElement(set_ctn_inner, qn('p:stCondLst'))
    set_cond = etree.SubElement(set_st, qn('p:cond'))
    set_cond.set('delay', '0')

    set_tgt = etree.SubElement(set_ctn, qn('p:tgtEl'))
    set_sp = etree.SubElement(set_tgt, qn('p:spTgt'))
    set_sp.set('spid', str(sp_id))

    attrNameLst = etree.SubElement(set_ctn, qn('p:attrNameLst'))
    attrName = etree.SubElement(attrNameLst, qn('p:attrName'))
    attrName.text = 'style.visibility'

    set_to = etree.SubElement(set_el, qn('p:to'))
    set_to_val = etree.SubElement(set_to, qn('p:strVal'))
    set_to_val.set('val', 'visible')


def add_slide_transition(slide, transition_type="fade", duration_ms=700):
    """Add slide transition."""
    sld = slide._element
    transition = sld.find(qn('p:transition'))
    if transition is None:
        transition = etree.SubElement(sld, qn('p:transition'))
    transition.set('spd', 'med')
    transition.set('advClick', '1')

    trans_map = {
        "fade":    qn('p:fade'),
        "push":    qn('p:push'),
        "wipe":    qn('p:wipe'),
        "cover":   qn('p:cover'),
        "split":   qn('p:split'),
    }
    if transition_type in trans_map:
        child = etree.SubElement(transition, trans_map[transition_type])
        if transition_type == "push":
            child.set('dir', 'l')


# ── Shape Helpers ──────────────────────────────────────────────

def add_bg(slide, color=COLORS["bg_dark"]):
    """Set dark background for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=COLORS["text_white"], bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=COLORS["text_light"], bullet_color=COLORS["accent_blue"]):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "▸ "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Segoe UI"
        # Item text
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = "Segoe UI"
    return txBox


def add_icon_label(slide, left, top, icon_text, label, icon_color, label_color=COLORS["text_white"],
                   icon_size=28, label_size=11):
    """Add an icon (emoji/symbol) with a label below."""
    icon_box = add_text_box(slide, left, top, Inches(1.2), Inches(0.5),
                            icon_text, font_size=icon_size, color=icon_color, alignment=PP_ALIGN.CENTER)
    label_box = add_text_box(slide, left, top + Inches(0.45), Inches(1.2), Inches(0.3),
                             label, font_size=label_size, color=label_color, alignment=PP_ALIGN.CENTER)
    return icon_box, label_box


def add_stage_pipeline(slide, left, top, stages, stage_color, width_per_stage=None):
    """Draw a horizontal pipeline of stages with arrows."""
    shapes = []
    n = len(stages)
    total_width = Inches(9.0)
    gap = Inches(0.08)
    stage_w = (total_width - gap * (n - 1)) / n
    if width_per_stage:
        stage_w = width_per_stage
    stage_h = Inches(0.45)

    for i, stage_name in enumerate(stages):
        x = left + i * (stage_w + gap)
        rect = add_rounded_rect(slide, int(x), top, int(stage_w), stage_h,
                                stage_color, border_color=COLORS["border"])
        # Add text to shape
        tf = rect.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = stage_name
        p.font.size = Pt(8)
        p.font.color.rgb = COLORS["white"]
        p.font.name = "Segoe UI Semibold"
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
        rect.text_frame.auto_size = None
        rect.text_frame.margin_top = Pt(4)
        rect.text_frame.margin_bottom = Pt(4)
        shapes.append(rect)

        # Arrow between stages
        if i < n - 1:
            arrow_x = int(x + stage_w + Emu(2000))
            arrow_box = add_text_box(slide, arrow_x, top, int(gap), stage_h,
                                     "›", font_size=14, color=COLORS["text_dim"],
                                     alignment=PP_ALIGN.CENTER)
    return shapes


def add_numbered_step(slide, left, top, number, title, description, accent_color):
    """Add a numbered step card."""
    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, Inches(0.4), Inches(0.4)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS["white"]
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    tf.margin_top = Pt(2)

    # Title
    title_box = add_text_box(slide, left + Inches(0.55), top + Inches(0.02),
                             Inches(3.0), Inches(0.3), title, font_size=14,
                             color=COLORS["text_white"], bold=True)
    # Description
    desc_box = add_text_box(slide, left + Inches(0.55), top + Inches(0.28),
                            Inches(3.5), Inches(0.3), description, font_size=10,
                            color=COLORS["text_light"])
    return circle, title_box, desc_box


def add_metric_card(slide, left, top, value, label, color):
    """Add a metric highlight card."""
    card = add_rounded_rect(slide, left, top, Inches(1.8), Inches(1.0),
                            COLORS["bg_card"], border_color=COLORS["border"])
    val_box = add_text_box(slide, left + Inches(0.1), top + Inches(0.1),
                           Inches(1.6), Inches(0.5), value, font_size=28,
                           color=color, bold=True, alignment=PP_ALIGN.CENTER)
    lbl_box = add_text_box(slide, left + Inches(0.1), top + Inches(0.55),
                           Inches(1.6), Inches(0.35), label, font_size=10,
                           color=COLORS["text_dim"], alignment=PP_ALIGN.CENTER)
    return card, val_box, lbl_box


# ══════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════

def build_title_slide(prs):
    """Slide 1: Title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide)
    add_slide_transition(slide, "fade")

    # Decorative top line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(1.5), Inches(2.6), Inches(7), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["accent_blue"]
    line.line.fill.background()

    # Title
    title = add_text_box(slide, Inches(1.5), Inches(2.75), Inches(7), Inches(0.8),
                         "AI-Powered DevOps Pipeline Platform",
                         font_size=32, color=COLORS["text_white"], bold=True,
                         alignment=PP_ALIGN.CENTER)
    add_fade_in(slide, title, delay_ms=200, duration_ms=800)

    # Subtitle
    subtitle = add_text_box(slide, Inches(1.5), Inches(3.55), Inches(7), Inches(0.5),
                            "Intelligent CI/CD Pipeline Generation with RAG + LLM",
                            font_size=16, color=COLORS["text_light"],
                            alignment=PP_ALIGN.CENTER)
    add_fade_in(slide, subtitle, delay_ms=600, duration_ms=600)

    # Three platform badges
    platforms = [
        ("GitLab CI", COLORS["gitlab_orange"]),
        ("Jenkins", COLORS["jenkins_red"]),
        ("GitHub Actions", COLORS["github_blue"]),
    ]
    badge_start = Inches(2.5)
    for i, (name, color) in enumerate(platforms):
        x = badge_start + i * Inches(1.8)
        badge = add_rounded_rect(slide, x, Inches(4.3), Inches(1.6), Inches(0.4),
                                 COLORS["bg_card"], border_color=color)
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(11)
        p.font.color.rgb = color
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
        tf.margin_top = Pt(3)
        add_fade_in(slide, badge, delay_ms=900 + i * 200, duration_ms=500)

    # Bottom text
    add_text_box(slide, Inches(1.5), Inches(5.2), Inches(7), Inches(0.3),
                 "Demo Presentation  ·  February 2026",
                 font_size=10, color=COLORS["text_dim"], alignment=PP_ALIGN.CENTER)


def build_overview_slide(prs):
    """Slide 2: Platform Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_transition(slide, "push")

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(5), Inches(0.5),
                 "Platform Overview", font_size=28, color=COLORS["text_white"], bold=True)
    add_text_box(slide, Inches(0.5), Inches(0.8), Inches(8), Inches(0.4),
                 "End-to-end intelligent pipeline generation across three CI/CD platforms",
                 font_size=12, color=COLORS["text_light"])

    # Separator
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.5), Inches(1.2), Inches(9), Pt(1))
    sep.fill.solid()
    sep.fill.fore_color.rgb = COLORS["border"]
    sep.line.fill.background()

    # Three column cards
    platforms = [
        {
            "name": "GitLab CI/CD",
            "color": COLORS["gitlab_orange"],
            "icon": "GL",
            "bullets": [
                "9-stage .gitlab-ci.yml generation",
                "GitLab server integration",
                "Artifact & cache management",
                "Docker-in-Docker pipelines",
            ],
            "repos": "13 repos  ·  13 templates"
        },
        {
            "name": "Jenkins",
            "color": COLORS["jenkins_red"],
            "icon": "JK",
            "bullets": [
                "9-stage Jenkinsfile generation",
                "Multibranch pipeline support",
                "Gitea SCM integration",
                "Docker agent orchestration",
            ],
            "repos": "10 repos  ·  10 templates"
        },
        {
            "name": "GitHub Actions",
            "color": COLORS["github_blue"],
            "icon": "GA",
            "bullets": [
                "10-job workflow YAML generation",
                "Gitea Actions runner compatible",
                "Multi-stage Dockerfile builds",
                "Shell-based (no marketplace actions)",
            ],
            "repos": "10 repos  ·  10 templates"
        },
    ]

    for i, plat in enumerate(platforms):
        x = Inches(0.4) + i * Inches(3.2)
        card = add_rounded_rect(slide, x, Inches(1.5), Inches(3.0), Inches(3.8),
                                COLORS["bg_card"], border_color=plat["color"])

        # Icon circle
        icon_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.15), Inches(1.7), Inches(0.45), Inches(0.45)
        )
        icon_circle.fill.solid()
        icon_circle.fill.fore_color.rgb = plat["color"]
        icon_circle.line.fill.background()
        tf = icon_circle.text_frame
        p = tf.paragraphs[0]
        p.text = plat["icon"]
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS["white"]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.margin_top = Pt(4)

        # Platform name
        add_text_box(slide, x + Inches(0.7), Inches(1.72), Inches(2.1), Inches(0.35),
                     plat["name"], font_size=16, color=plat["color"], bold=True)

        # Bullet points
        add_bullet_list(slide, x + Inches(0.15), Inches(2.25), Inches(2.7), Inches(2.0),
                        plat["bullets"], font_size=11, bullet_color=plat["color"])

        # Repo count badge
        badge = add_rounded_rect(slide, x + Inches(0.2), Inches(4.7), Inches(2.6), Inches(0.35),
                                 COLORS["bg_dark"], border_color=COLORS["border"])
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = plat["repos"]
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS["text_dim"]
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
        tf.margin_top = Pt(3)

        add_fade_in(slide, card, delay_ms=200 + i * 300, duration_ms=600)


def build_architecture_slide(prs):
    """Slide 3: Architecture & Tech Stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_transition(slide, "fade")

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(5), Inches(0.5),
                 "Architecture", font_size=28, color=COLORS["text_white"], bold=True)

    # ── Flow diagram: User → Chat UI → Backend → LLM/RAG → CI Platform ──
    flow_y = Inches(1.3)
    flow_boxes = [
        ("User\nChat UI", COLORS["accent_blue"]),
        ("FastAPI\nBackend", COLORS["accent_green"]),
        ("RAG\nChromaDB", COLORS["accent_purple"]),
        ("LLM\nOllama/Claude/OpenAI", COLORS["accent_orange"]),
        ("CI Platform\nGitLab/Jenkins/Gitea", COLORS["gitlab_orange"]),
    ]
    box_w = Inches(1.6)
    box_h = Inches(0.7)
    gap = Inches(0.25)
    start_x = Inches(0.4)

    for i, (label, color) in enumerate(flow_boxes):
        x = start_x + i * (box_w + gap)
        rect = add_rounded_rect(slide, int(x), flow_y, int(box_w), box_h,
                                COLORS["bg_card"], border_color=color)
        tf = rect.text_frame
        tf.word_wrap = True
        for li, line in enumerate(label.split('\n')):
            if li == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(9) if li == 0 else Pt(8)
            p.font.color.rgb = color if li == 0 else COLORS["text_light"]
            p.font.bold = (li == 0)
            p.font.name = "Segoe UI"
            p.alignment = PP_ALIGN.CENTER
        tf.margin_top = Pt(6)
        add_fade_in(slide, rect, delay_ms=200 + i * 200, duration_ms=500)

        if i < len(flow_boxes) - 1:
            ax = int(x + box_w + Emu(5000))
            arr = add_text_box(slide, ax, flow_y + Inches(0.15), int(gap - Emu(10000)), Inches(0.4),
                               "→", font_size=18, color=COLORS["text_dim"], alignment=PP_ALIGN.CENTER)

    # ── Tech Stack cards ──
    stack_y = Inches(2.4)
    add_text_box(slide, Inches(0.5), stack_y, Inches(3), Inches(0.35),
                 "Tech Stack", font_size=16, color=COLORS["accent_blue"], bold=True)

    tech_items = [
        ("Backend",    "FastAPI · Python 3.11 · httpx · Pydantic v2"),
        ("LLM",        "Ollama (qwen3:32b) · Claude (Opus/Sonnet) · OpenAI (GPT-4)"),
        ("RAG Store",  "ChromaDB · Sentence Transformers · Vector similarity"),
        ("Git Servers", "GitLab CE · Gitea (Jenkins + GitHub Actions)"),
        ("Registry",   "Nexus 3 OSS · Docker Registry (port 5001)"),
        ("Security",   "SonarQube · Trivy · Semgrep (SAST)"),
        ("Secrets",    "HashiCorp Vault · KV-v2 · AppRole auth"),
        ("Infra",      "Docker Compose · 31 containers · 5 networks"),
    ]

    for i, (label, desc) in enumerate(tech_items):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(4.8)
        y = stack_y + Inches(0.4) + row * Inches(0.55)

        lbl = add_text_box(slide, x, y, Inches(1.1), Inches(0.3),
                           label, font_size=10, color=COLORS["accent_green"], bold=True)
        val = add_text_box(slide, x + Inches(1.1), y, Inches(3.5), Inches(0.3),
                           desc, font_size=10, color=COLORS["text_light"])

    # ── Metrics row ──
    metrics_y = Inches(5.0)
    metrics = [
        ("25", "Platform Tools", COLORS["accent_blue"]),
        ("33", "RAG Templates", COLORS["accent_purple"]),
        ("18+", "Languages", COLORS["accent_green"]),
        ("31", "Docker Services", COLORS["accent_orange"]),
        ("3", "LLM Providers", COLORS["gitlab_orange"]),
    ]
    for i, (val, lbl, clr) in enumerate(metrics):
        x = Inches(0.3) + i * Inches(1.9)
        c, v, l = add_metric_card(slide, x, metrics_y, val, lbl, clr)
        add_fade_in(slide, c, delay_ms=800 + i * 150, duration_ms=400)


def build_demo_flow_slide(prs):
    """Slide 4: How the Demo Works (conversation flow)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_transition(slide, "push")

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(6), Inches(0.5),
                 "Chat-Driven Pipeline Generation", font_size=28,
                 color=COLORS["text_white"], bold=True)
    add_text_box(slide, Inches(0.5), Inches(0.8), Inches(8), Inches(0.35),
                 "Natural language conversation → production-ready CI/CD pipeline",
                 font_size=12, color=COLORS["text_light"])

    # Steps
    steps = [
        ("Provide Repo URL",
         "User pastes a repository URL in the chat interface",
         COLORS["accent_blue"]),
        ("Analyze Repository",
         "Backend clones repo, detects language, framework, build tool, dependencies",
         COLORS["accent_purple"]),
        ("Query RAG Templates",
         "ChromaDB searched for matching successful pipeline templates",
         COLORS["accent_green"]),
        ("Generate Pipeline",
         "LLM generates .gitlab-ci.yml / Jenkinsfile / workflow YAML + Dockerfile",
         COLORS["accent_orange"]),
        ("Validate & Fix",
         "Automated YAML validation, image verification, iterative LLM self-healing",
         COLORS["accent_red"]),
        ("Review & Approve",
         "User reviews generated config, approves or requests changes",
         COLORS["text_white"]),
        ("Commit & Push",
         "Pipeline committed to feature branch via Git API",
         COLORS["accent_blue"]),
        ("Monitor Build",
         "Real-time build status tracking, auto-triggered by commit",
         COLORS["accent_purple"]),
        ("Learn on Success",
         "Successful pipeline stored in ChromaDB for future RAG retrieval",
         COLORS["accent_green"]),
    ]

    for i, (title, desc, color) in enumerate(steps):
        col = i // 5
        row = i % 5
        x = Inches(0.4) + col * Inches(4.8)
        y = Inches(1.35) + row * Inches(0.95)
        c, t, d = add_numbered_step(slide, x, y, i + 1, title, desc, color)
        add_fade_in(slide, c, delay_ms=200 + i * 200, duration_ms=400)
        add_fade_in(slide, t, delay_ms=300 + i * 200, duration_ms=400)
        add_fade_in(slide, d, delay_ms=350 + i * 200, duration_ms=400)


def build_gitlab_slide(prs):
    """Slide 5: GitLab Pipeline Generator Deep Dive."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_transition(slide, "push")

    # Header with brand color
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0), Inches(0), Inches(10), Inches(0.08))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLORS["gitlab_orange"]
    header_bar.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(6), Inches(0.5),
                 "GitLab CI/CD Pipeline Generator", font_size=26,
                 color=COLORS["gitlab_orange"], bold=True)
    add_text_box(slide, Inches(0.5), Inches(0.7), Inches(8), Inches(0.35),
                 "Generates .gitlab-ci.yml + Dockerfile for GitLab CE with Docker-in-Docker runner",
                 font_size=11, color=COLORS["text_light"])

    # Pipeline stages
    stages = ["Compile", "Build", "Test", "SAST", "SonarQube",
              "Trivy Scan", "Push", "Notify", "Learn"]
    stage_shapes = add_stage_pipeline(slide, Inches(0.5), Inches(1.2), stages,
                                      COLORS["gitlab_orange"])
    for i, s in enumerate(stage_shapes):
        add_fade_in(slide, s, delay_ms=300 + i * 100, duration_ms=300)

    # Left: Features
    feat_y = Inches(1.9)
    add_text_box(slide, Inches(0.5), feat_y, Inches(4), Inches(0.3),
                 "Key Features", font_size=14, color=COLORS["gitlab_orange"], bold=True)
    features = [
        "Language auto-detection (Java, Python, Go, Node.js, Ruby, Rust, .NET)",
        "Framework-aware generation (Spring Boot, FastAPI, Gin, Express, etc.)",
        "DinD network isolation — pre-built Nexus images, no apk/apt",
        "RAG priority: language+framework+build_tool match first",
        "Iterative LLM self-healing with error pattern matching",
        "Automatic pipeline learning on successful builds",
        "Splunk + Jira notifications on failure",
    ]
    feat_box = add_bullet_list(slide, Inches(0.5), feat_y + Inches(0.3), Inches(4.5), Inches(3.0),
                               features, font_size=10, bullet_color=COLORS["gitlab_orange"])
    add_fade_in(slide, feat_box, delay_ms=800, duration_ms=500)

    # Right: Chat example
    chat_y = Inches(1.9)
    add_text_box(slide, Inches(5.3), chat_y, Inches(4), Inches(0.3),
                 "Chat Interaction", font_size=14, color=COLORS["gitlab_orange"], bold=True)

    chat_msgs = [
        ("User", "Generate a pipeline for\nhttp://gitlab-server/ai-pipeline-\nprojects/java-springboot-api",
         COLORS["accent_blue"], COLORS["bg_card"]),
        ("AI", "Analyzing repository...\nDetected: Java · Spring Boot · Maven\nQuerying RAG templates...",
         COLORS["accent_green"], COLORS["bg_card"]),
        ("AI", "Pipeline generated! 9 stages\nSource: RAG template match\n[Approve] [Modify] [Reject]",
         COLORS["accent_green"], COLORS["bg_card"]),
        ("User", "Approve",
         COLORS["accent_blue"], COLORS["bg_card"]),
        ("AI", "Committed to feature branch.\nBuild running... all 9 stages passed ✓\nTemplate saved to ChromaDB.",
         COLORS["accent_green"], COLORS["bg_card"]),
    ]

    for i, (sender, msg, sender_color, bg) in enumerate(chat_msgs):
        y = chat_y + Inches(0.35) + i * Inches(0.72)
        bubble = add_rounded_rect(slide, Inches(5.3), y, Inches(4.2), Inches(0.65),
                                  bg, border_color=COLORS["border"])
        add_text_box(slide, Inches(5.45), y + Inches(0.03), Inches(0.6), Inches(0.2),
                     sender, font_size=8, color=sender_color, bold=True)
        add_text_box(slide, Inches(5.45), y + Inches(0.2), Inches(3.9), Inches(0.42),
                     msg, font_size=8, color=COLORS["text_light"])
        add_fade_in(slide, bubble, delay_ms=1000 + i * 300, duration_ms=400)


def build_jenkins_slide(prs):
    """Slide 6: Jenkins Pipeline Generator Deep Dive."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_transition(slide, "push")

    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0), Inches(0), Inches(10), Inches(0.08))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLORS["jenkins_red"]
    header_bar.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(6), Inches(0.5),
                 "Jenkins Pipeline Generator", font_size=26,
                 color=COLORS["jenkins_red"], bold=True)
    add_text_box(slide, Inches(0.5), Inches(0.7), Inches(8), Inches(0.35),
                 "Generates Jenkinsfile + Dockerfile for declarative multibranch pipelines with Docker agents",
                 font_size=11, color=COLORS["text_light"])

    # Pipeline stages
    stages = ["Compile", "Build Image", "Test Image", "Static Analysis",
              "SonarQube", "Trivy Scan", "Push Release", "Notify", "Learn"]
    stage_shapes = add_stage_pipeline(slide, Inches(0.5), Inches(1.2), stages,
                                      COLORS["jenkins_red"])
    for i, s in enumerate(stage_shapes):
        add_fade_in(slide, s, delay_ms=300 + i * 100, duration_ms=300)

    # Left column
    feat_y = Inches(1.9)
    add_text_box(slide, Inches(0.5), feat_y, Inches(4), Inches(0.3),
                 "Key Features", font_size=14, color=COLORS["jenkins_red"], bold=True)
    features = [
        "Declarative Jenkinsfile with Docker agent label",
        "Gitea SCM integration (not GitLab)",
        "Multibranch pipeline auto-discovery",
        "Credential management via Jenkins credential store",
        "Pattern-based URL detection in chat",
        "Background build monitoring with auto-learning",
        "Branch scan trigger on new commits",
    ]
    feat_box = add_bullet_list(slide, Inches(0.5), feat_y + Inches(0.3), Inches(4.5), Inches(3.0),
                               features, font_size=10, bullet_color=COLORS["jenkins_red"])
    add_fade_in(slide, feat_box, delay_ms=800, duration_ms=500)

    # Right: Architecture detail
    arch_y = Inches(1.9)
    add_text_box(slide, Inches(5.3), arch_y, Inches(4), Inches(0.3),
                 "Jenkins Architecture", font_size=14, color=COLORS["jenkins_red"], bold=True)

    arch_items = [
        ("Git Server", "Gitea (gitea-server:3000)"),
        ("Organization", "jenkins-projects"),
        ("Build Agents", "3 Docker agents with docker label"),
        ("Registry", "Nexus at localhost:5001"),
        ("Monitoring", "Real-time build status polling"),
        ("Context Path", "/jenkins (JENKINS_OPTS)"),
    ]
    for i, (k, v) in enumerate(arch_items):
        y = arch_y + Inches(0.4) + i * Inches(0.48)
        card = add_rounded_rect(slide, Inches(5.3), y, Inches(4.2), Inches(0.4),
                                COLORS["bg_card"], border_color=COLORS["border"])
        add_text_box(slide, Inches(5.45), y + Inches(0.05), Inches(1.3), Inches(0.3),
                     k, font_size=9, color=COLORS["jenkins_red"], bold=True)
        add_text_box(slide, Inches(6.8), y + Inches(0.05), Inches(2.6), Inches(0.3),
                     v, font_size=9, color=COLORS["text_light"])
        add_fade_in(slide, card, delay_ms=1000 + i * 200, duration_ms=400)


def build_github_slide(prs):
    """Slide 7: GitHub Actions Pipeline Generator Deep Dive."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_transition(slide, "push")

    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0), Inches(0), Inches(10), Inches(0.08))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLORS["github_blue"]
    header_bar.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(6), Inches(0.5),
                 "GitHub Actions Pipeline Generator", font_size=26,
                 color=COLORS["github_blue"], bold=True)
    add_text_box(slide, Inches(0.5), Inches(0.7), Inches(8), Inches(0.35),
                 "Generates workflow YAML + Dockerfile for Gitea Actions runner (GitHub Actions compatible)",
                 font_size=11, color=COLORS["text_light"])

    # Pipeline stages (10 jobs)
    stages = ["Compile", "Build", "Test", "Static\nAnalysis", "SonarQube",
              "Trivy", "Push", "Notify OK", "Notify Fail", "Learn"]
    stage_shapes = add_stage_pipeline(slide, Inches(0.5), Inches(1.2), stages,
                                      COLORS["github_blue"])
    for i, s in enumerate(stage_shapes):
        add_fade_in(slide, s, delay_ms=300 + i * 100, duration_ms=300)

    # Left: Features
    feat_y = Inches(1.9)
    add_text_box(slide, Inches(0.5), feat_y, Inches(4), Inches(0.3),
                 "Key Features", font_size=14, color=COLORS["github_blue"], bold=True)
    features = [
        "Gitea Actions runner (Alpine-based act_runner)",
        "Shell-based checkout (no actions/checkout — cache corruption fix)",
        "Token-authenticated git clone for private repos",
        "No Docker marketplace actions — shell docker build/push",
        "Multi-stage Dockerfiles (no artifact passing support)",
        "wget instead of curl (Alpine runner limitation)",
        "Default templates preferred over LLM for known languages",
        "E2E verified: 10 repos × 10 jobs all passing",
    ]
    feat_box = add_bullet_list(slide, Inches(0.5), feat_y + Inches(0.3), Inches(4.5), Inches(3.2),
                               features, font_size=10, bullet_color=COLORS["github_blue"])
    add_fade_in(slide, feat_box, delay_ms=800, duration_ms=500)

    # Right: Gitea quirks
    quirk_y = Inches(1.9)
    add_text_box(slide, Inches(5.3), quirk_y, Inches(4), Inches(0.3),
                 "Gitea Actions Adaptations", font_size=14, color=COLORS["github_blue"], bold=True)

    quirks = [
        ("No Artifacts", "Multi-stage Docker builds instead of\nupload-artifact/download-artifact"),
        ("No Docker Actions", "Shell docker build/push replaces\ndocker/build-push-action"),
        ("Runner Cache", "Shell git clone avoids stale\nactions/checkout v4 cache"),
        ("Docker Networking", "localhost:5001 for registry\n(host daemon, not container DNS)"),
        ("YAML on: fix", "Post-process yaml.dump to restore\n'on:' key from boolean 'true:'"),
    ]
    for i, (title, desc) in enumerate(quirks):
        y = quirk_y + Inches(0.4) + i * Inches(0.6)
        card = add_rounded_rect(slide, Inches(5.3), y, Inches(4.2), Inches(0.52),
                                COLORS["bg_card"], border_color=COLORS["border"])
        add_text_box(slide, Inches(5.45), y + Inches(0.02), Inches(1.5), Inches(0.22),
                     title, font_size=9, color=COLORS["github_blue"], bold=True)
        add_text_box(slide, Inches(6.95), y + Inches(0.02), Inches(2.4), Inches(0.48),
                     desc, font_size=8, color=COLORS["text_light"])
        add_fade_in(slide, card, delay_ms=1000 + i * 200, duration_ms=400)


def build_rag_learning_slide(prs):
    """Slide 8: RAG Learning Loop."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_transition(slide, "fade")

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(7), Inches(0.5),
                 "RAG Learning Loop — Self-Improving Pipelines", font_size=26,
                 color=COLORS["accent_purple"], bold=True)
    add_text_box(slide, Inches(0.5), Inches(0.8), Inches(8), Inches(0.35),
                 "Every successful pipeline build becomes a template for future generations",
                 font_size=12, color=COLORS["text_light"])

    # Circular flow
    loop_steps = [
        ("1. Generate", "LLM creates pipeline config\nfrom repo analysis + RAG match", COLORS["accent_orange"], Inches(0.4), Inches(1.6)),
        ("2. Validate", "YAML syntax, image availability,\nstage structure verification", COLORS["accent_red"], Inches(3.5), Inches(1.6)),
        ("3. Commit", "Push to feature branch\nvia Git API (GitLab/Gitea)", COLORS["accent_blue"], Inches(6.6), Inches(1.6)),
        ("4. Build", "CI platform executes pipeline\nReal-time status monitoring", COLORS["accent_green"], Inches(6.6), Inches(3.2)),
        ("5. Learn", "On success → store in ChromaDB\nwith metadata + embeddings", COLORS["accent_purple"], Inches(3.5), Inches(3.2)),
        ("6. Retrieve", "Next generation queries ChromaDB\nfor best-match template", COLORS["accent_blue"], Inches(0.4), Inches(3.2)),
    ]

    for i, (title, desc, color, x, y) in enumerate(loop_steps):
        card = add_rounded_rect(slide, x, y, Inches(2.8), Inches(1.1),
                                COLORS["bg_card"], border_color=color)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.08), Inches(2.6), Inches(0.3),
                     title, font_size=13, color=color, bold=True)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.4), Inches(2.6), Inches(0.65),
                     desc, font_size=10, color=COLORS["text_light"])
        add_fade_in(slide, card, delay_ms=300 + i * 300, duration_ms=500)

    # Arrows between steps (horizontal and vertical connectors)
    arrows = ["→", "→", "↓", "←", "←", "↑"]
    arrow_positions = [
        (Inches(3.1), Inches(1.95), Inches(0.5), Inches(0.3)),
        (Inches(6.2), Inches(1.95), Inches(0.5), Inches(0.3)),
        (Inches(7.5), Inches(2.7), Inches(0.5), Inches(0.5)),
        (Inches(6.2), Inches(3.55), Inches(0.5), Inches(0.3)),
        (Inches(3.1), Inches(3.55), Inches(0.5), Inches(0.3)),
        (Inches(1.3), Inches(2.7), Inches(0.5), Inches(0.5)),
    ]
    for arrow, (ax, ay, aw, ah) in zip(arrows, arrow_positions):
        add_text_box(slide, ax, ay, aw, ah, arrow, font_size=18,
                     color=COLORS["text_dim"], alignment=PP_ALIGN.CENTER)

    # ChromaDB detail card
    detail_y = Inches(4.7)
    detail_card = add_rounded_rect(slide, Inches(0.4), detail_y, Inches(9.2), Inches(0.85),
                                   COLORS["bg_card"], border_color=COLORS["accent_purple"])
    add_text_box(slide, Inches(0.6), detail_y + Inches(0.08), Inches(2), Inches(0.25),
                 "ChromaDB Collections", font_size=11, color=COLORS["accent_purple"], bold=True)

    collections = [
        ("gitlab_successful_pipelines", "13 templates"),
        ("jenkins_successful_pipelines", "10 templates"),
        ("github_actions_successful_pipelines", "10 templates"),
    ]
    for i, (name, count) in enumerate(collections):
        x = Inches(0.6) + i * Inches(3.1)
        add_text_box(slide, x, detail_y + Inches(0.38), Inches(2.8), Inches(0.2),
                     name, font_size=8, color=COLORS["text_light"], font_name="Consolas")
        add_text_box(slide, x, detail_y + Inches(0.55), Inches(1), Inches(0.2),
                     count, font_size=8, color=COLORS["accent_green"], bold=True)
    add_fade_in(slide, detail_card, delay_ms=2200, duration_ms=500)


def build_validation_slide(prs):
    """Slide 9: Validation & Self-Healing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_transition(slide, "push")

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(7), Inches(0.5),
                 "Validation & LLM Self-Healing", font_size=26,
                 color=COLORS["accent_green"], bold=True)
    add_text_box(slide, Inches(0.5), Inches(0.8), Inches(8), Inches(0.35),
                 "Multi-layer validation with automatic error correction before any pipeline reaches CI",
                 font_size=12, color=COLORS["text_light"])

    # Validation layers
    layers = [
        ("Layer 1: Static Validation", COLORS["accent_blue"], [
            "YAML/Groovy syntax parsing",
            "Required stages verification (compile → learn)",
            "Docker image reference validation",
            "Environment variable completeness check",
        ]),
        ("Layer 2: Image Verification", COLORS["accent_orange"], [
            "Public-to-Nexus image mapping",
            "Registry connectivity test",
            "Base image availability check",
            "Version compatibility (e.g. rust:1.93-slim)",
        ]),
        ("Layer 3: LLM Self-Healing", COLORS["accent_purple"], [
            "Error pattern detection (40+ patterns)",
            "LLM receives error + pipeline + context",
            "Iterative fix (up to 3 attempts)",
            "Patterns: TLS errors, missing deps, syntax, artifact size",
        ]),
    ]

    for i, (title, color, items) in enumerate(layers):
        y = Inches(1.4) + i * Inches(1.6)
        card = add_rounded_rect(slide, Inches(0.4), y, Inches(9.2), Inches(1.3),
                                COLORS["bg_card"], border_color=color)
        add_text_box(slide, Inches(0.6), y + Inches(0.08), Inches(3), Inches(0.3),
                     title, font_size=13, color=color, bold=True)

        for j, item in enumerate(items):
            col = j // 2
            row = j % 2
            x = Inches(0.6) + col * Inches(4.5)
            iy = y + Inches(0.42) + row * Inches(0.35)
            add_text_box(slide, x, iy, Inches(4.3), Inches(0.3),
                         f"▸ {item}", font_size=10, color=COLORS["text_light"])

        add_fade_in(slide, card, delay_ms=300 + i * 400, duration_ms=600)


def build_live_demo_slide(prs):
    """Slide 10: Live Demo slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_transition(slide, "fade")

    # Center content
    title = add_text_box(slide, Inches(1.5), Inches(1.8), Inches(7), Inches(0.8),
                         "Live Demo", font_size=42,
                         color=COLORS["text_white"], bold=True,
                         alignment=PP_ALIGN.CENTER)
    add_fade_in(slide, title, delay_ms=200, duration_ms=800)

    # Decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(3.5), Inches(2.7), Inches(3), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["accent_blue"]
    line.line.fill.background()

    subtitle = add_text_box(slide, Inches(1.5), Inches(2.9), Inches(7), Inches(0.5),
                            "Generating pipelines across all three platforms",
                            font_size=16, color=COLORS["text_light"],
                            alignment=PP_ALIGN.CENTER)
    add_fade_in(slide, subtitle, delay_ms=600, duration_ms=600)

    # Demo steps
    demo_items = [
        ("1", "GitLab CI/CD", "Java Spring Boot API → .gitlab-ci.yml", COLORS["gitlab_orange"]),
        ("2", "Jenkins", "Python FastAPI Service → Jenkinsfile", COLORS["jenkins_red"]),
        ("3", "GitHub Actions", "Go Gin API → workflow YAML", COLORS["github_blue"]),
    ]

    for i, (num, platform, desc, color) in enumerate(demo_items):
        y = Inches(3.6) + i * Inches(0.65)
        badge = add_rounded_rect(slide, Inches(2.0), y, Inches(6), Inches(0.5),
                                 COLORS["bg_card"], border_color=color)
        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(2.15), y + Inches(0.06), Inches(0.35), Inches(0.35)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS["white"]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.margin_top = Pt(2)

        add_text_box(slide, Inches(2.65), y + Inches(0.05), Inches(1.5), Inches(0.3),
                     platform, font_size=12, color=color, bold=True)
        add_text_box(slide, Inches(4.2), y + Inches(0.08), Inches(3.5), Inches(0.3),
                     desc, font_size=11, color=COLORS["text_light"])
        add_fade_in(slide, badge, delay_ms=800 + i * 300, duration_ms=500)


def build_comparison_slide(prs):
    """Slide 11: Platform Comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_transition(slide, "push")

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(7), Inches(0.5),
                 "Platform Comparison", font_size=26,
                 color=COLORS["text_white"], bold=True)

    # Table header
    headers = ["Feature", "GitLab CI", "Jenkins", "GitHub Actions"]
    header_colors = [COLORS["text_dim"], COLORS["gitlab_orange"], COLORS["jenkins_red"], COLORS["github_blue"]]
    col_widths = [Inches(2.8), Inches(2.0), Inches(2.0), Inches(2.0)]
    col_x = [Inches(0.5)]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w + Inches(0.15))

    for i, (header, color) in enumerate(zip(headers, header_colors)):
        add_text_box(slide, col_x[i], Inches(1.0), col_widths[i], Inches(0.35),
                     header, font_size=11, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    # Table rows
    rows = [
        ("Pipeline Format", ".gitlab-ci.yml", "Jenkinsfile", "workflow YAML"),
        ("Git Server", "GitLab CE", "Gitea", "Gitea"),
        ("Runner Type", "Docker-in-Docker", "Docker Agents (×3)", "act_runner (Alpine)"),
        ("Pipeline Stages", "9 stages", "9 stages", "10 jobs"),
        ("Registry", "Nexus (ai-nexus:5001)", "Nexus (localhost:5001)", "Nexus (localhost:5001)"),
        ("Checkout Method", "GitLab CI built-in", "SCM plugin", "Shell git clone"),
        ("Artifact Strategy", "GitLab artifacts", "Stash/unstash", "Multi-stage Dockerfile"),
        ("Chat Interface", "Ollama tool-calling", "Pattern detection", "Pattern detection"),
        ("LLM Self-Heal", "3 iterations max", "3 iterations max", "3 iterations max"),
        ("RAG Templates", "13 in ChromaDB", "10 in ChromaDB", "10 in ChromaDB"),
    ]

    for ri, (feature, *values) in enumerate(rows):
        y = Inches(1.45) + ri * Inches(0.42)
        bg_color = COLORS["bg_card"] if ri % 2 == 0 else COLORS["bg_dark"]
        row_bg = add_rounded_rect(slide, Inches(0.4), y, Inches(9.2), Inches(0.38),
                                  bg_color)
        add_text_box(slide, col_x[0], y + Inches(0.05), col_widths[0], Inches(0.28),
                     feature, font_size=9, color=COLORS["text_white"], bold=True)
        for vi, val in enumerate(values):
            add_text_box(slide, col_x[vi + 1], y + Inches(0.05), col_widths[vi + 1], Inches(0.28),
                         val, font_size=9, color=COLORS["text_light"],
                         alignment=PP_ALIGN.CENTER)


def build_closing_slide(prs):
    """Slide 12: Closing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_transition(slide, "fade")

    title = add_text_box(slide, Inches(1.5), Inches(1.5), Inches(7), Inches(0.8),
                         "Thank You", font_size=42,
                         color=COLORS["text_white"], bold=True,
                         alignment=PP_ALIGN.CENTER)
    add_fade_in(slide, title, delay_ms=200, duration_ms=800)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(3.5), Inches(2.4), Inches(3), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["accent_blue"]
    line.line.fill.background()

    subtitle = add_text_box(slide, Inches(1.0), Inches(2.7), Inches(8), Inches(0.5),
                            "AI-Powered DevOps Pipeline Platform",
                            font_size=18, color=COLORS["text_light"],
                            alignment=PP_ALIGN.CENTER)
    add_fade_in(slide, subtitle, delay_ms=600, duration_ms=600)

    # Summary stats
    stats = [
        ("3", "CI/CD\nPlatforms", COLORS["accent_blue"]),
        ("33+", "Managed\nRepositories", COLORS["accent_green"]),
        ("33", "RAG\nTemplates", COLORS["accent_purple"]),
        ("9-10", "Pipeline\nStages", COLORS["accent_orange"]),
    ]
    for i, (val, lbl, color) in enumerate(stats):
        x = Inches(1.2) + i * Inches(2.0)
        y = Inches(3.5)
        card = add_rounded_rect(slide, x, y, Inches(1.7), Inches(1.2),
                                COLORS["bg_card"], border_color=color)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.1), Inches(1.5), Inches(0.5),
                     val, font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.6), Inches(1.5), Inches(0.5),
                     lbl, font_size=10, color=COLORS["text_dim"], alignment=PP_ALIGN.CENTER)
        add_fade_in(slide, card, delay_ms=800 + i * 200, duration_ms=500)

    # URLs
    urls = [
        "Backend: http://localhost:8003",
        "GitLab: http://localhost:8929",
        "Jenkins: http://localhost:8080/jenkins",
        "Gitea: http://localhost:3002",
    ]
    for i, url in enumerate(urls):
        add_text_box(slide, Inches(2.0), Inches(5.0) + i * Inches(0.25), Inches(6), Inches(0.25),
                     url, font_size=9, color=COLORS["text_dim"],
                     alignment=PP_ALIGN.CENTER, font_name="Consolas")


def build_frontend_slide(prs):
    """Slide: Frontend Dashboard — 25 Tools, 5 Categories."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_transition(slide, "push")

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(7), Inches(0.5),
                 "Unified Dashboard — 25 Tools", font_size=26,
                 color=COLORS["text_white"], bold=True)
    add_text_box(slide, Inches(0.5), Inches(0.75), Inches(8), Inches(0.35),
                 "Single-page application with chat-driven interfaces and real-time progress tracking",
                 font_size=12, color=COLORS["text_light"])

    categories = [
        {
            "name": "DevOps",
            "color": COLORS["accent_blue"],
            "tools": ["GitLab Pipeline Gen", "Jenkins Pipeline Gen",
                      "GitHub Actions Gen", "AWS CodePipeline*", "Azure DevOps*"],
        },
        {
            "name": "IaC",
            "color": COLORS["accent_green"],
            "tools": ["Terraform Generator", "Ansible Playbooks*",
                      "K8s Manifests*", "CloudFormation*"],
        },
        {
            "name": "Support",
            "color": COLORS["accent_orange"],
            "tools": ["Incident Assistant*", "Log Analyzer*",
                      "Troubleshooting*", "Knowledge Base*"],
        },
        {
            "name": "SRE",
            "color": COLORS["accent_red"],
            "tools": ["SLO Calculator*", "Capacity Planner*",
                      "Runbook Generator*", "Postmortem Gen*"],
        },
        {
            "name": "Auxiliary",
            "color": COLORS["accent_purple"],
            "tools": ["Connectivity Validator", "Commit History",
                      "ChromaDB Browser", "Secret Manager",
                      "Dependency Scanner", "Release Notes",
                      "Migration Assistant", "Compliance Check"],
        },
    ]

    for ci, cat in enumerate(categories):
        col = ci % 3
        row = ci // 3
        x = Inches(0.3) + col * Inches(3.2)
        y = Inches(1.3) + row * Inches(2.0)
        card_h = Inches(1.75) if ci < 3 else Inches(1.75)

        card = add_rounded_rect(slide, x, y, Inches(3.05), card_h,
                                COLORS["bg_card"], border_color=cat["color"])

        # Category header
        add_text_box(slide, x + Inches(0.1), y + Inches(0.08), Inches(2.8), Inches(0.28),
                     cat["name"], font_size=13, color=cat["color"], bold=True)

        # Tools as mini-badges
        for ti, tool in enumerate(cat["tools"]):
            tcol = ti % 2
            trow = ti // 2
            tx = x + Inches(0.1) + tcol * Inches(1.45)
            ty = y + Inches(0.4) + trow * Inches(0.3)
            is_upcoming = tool.endswith("*")
            display_name = tool.rstrip("*")
            text_color = COLORS["text_dim"] if is_upcoming else COLORS["text_light"]
            add_text_box(slide, tx, ty, Inches(1.4), Inches(0.25),
                         display_name, font_size=8, color=text_color)

        add_fade_in(slide, card, delay_ms=200 + ci * 250, duration_ms=500)

    # Legend
    add_text_box(slide, Inches(0.5), Inches(5.15), Inches(5), Inches(0.25),
                 "* Coming soon    |    Active tools include chat UI, progress tracking, LLM switching",
                 font_size=9, color=COLORS["text_dim"])


def build_languages_slide(prs):
    """Slide: Supported Languages & Frameworks."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_transition(slide, "fade")

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(7), Inches(0.5),
                 "18+ Languages & Frameworks", font_size=26,
                 color=COLORS["accent_green"], bold=True)
    add_text_box(slide, Inches(0.5), Inches(0.75), Inches(8), Inches(0.35),
                 "Auto-detected from repository analysis with pre-built Nexus images for air-gapped DinD execution",
                 font_size=12, color=COLORS["text_light"])

    languages = [
        ("Java",       "Maven / Gradle",          "Spring Boot, Quarkus",    "maven:3.9-eclipse-temurin-17"),
        ("Kotlin",     "Gradle",                   "Ktor, Spring",           "gradle:8.7-jdk17-alpine"),
        ("Scala",      "sbt",                      "Play, Akka",            "sbt:1.9-eclipse-temurin-17"),
        ("Python",     "pip / poetry",             "FastAPI, Django, Flask", "python:3.11-slim"),
        ("JavaScript", "npm / yarn",               "Express, React, Next",  "node:20-alpine"),
        ("TypeScript", "npm / yarn",               "NestJS, Deno",          "node:20-alpine"),
        ("Go",         "go mod",                   "Gin, Echo, Fiber",      "golang:1.22-alpine-git"),
        ("Rust",       "cargo",                    "Actix, Axum, Rocket",   "rust:1.93-slim"),
        ("Ruby",       "bundler / gem",            "Rails, Sinatra",        "ruby:3.3-alpine"),
        ("PHP",        "composer",                 "Laravel, Symfony",      "php:8.3-fpm-alpine"),
        ("C# / .NET",  "dotnet",                   "ASP.NET, Blazor",      "dotnet-aspnet:8.0-alpine"),
        ("Perl",       "cpanm",                    "Mojolicious, Dancer",   "perl:5.38-slim"),
    ]

    # Table header
    headers = ["Language", "Build Tool", "Frameworks", "Base Image"]
    header_x = [Inches(0.4), Inches(1.8), Inches(3.4), Inches(6.2)]
    header_w = [Inches(1.3), Inches(1.5), Inches(2.7), Inches(3.2)]

    for i, h in enumerate(headers):
        add_text_box(slide, header_x[i], Inches(1.2), header_w[i], Inches(0.3),
                     h, font_size=10, color=COLORS["accent_green"], bold=True)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.4), Inches(1.48), Inches(9.2), Pt(1))
    sep.fill.solid()
    sep.fill.fore_color.rgb = COLORS["border"]
    sep.line.fill.background()

    for ri, (lang, build, frameworks, image) in enumerate(languages):
        y = Inches(1.55) + ri * Inches(0.33)
        bg_color = COLORS["bg_card"] if ri % 2 == 0 else COLORS["bg_dark"]
        row_bg = add_rounded_rect(slide, Inches(0.35), y, Inches(9.3), Inches(0.3), bg_color)
        add_text_box(slide, header_x[0], y + Inches(0.03), header_w[0], Inches(0.24),
                     lang, font_size=9, color=COLORS["text_white"], bold=True)
        add_text_box(slide, header_x[1], y + Inches(0.03), header_w[1], Inches(0.24),
                     build, font_size=9, color=COLORS["text_light"])
        add_text_box(slide, header_x[2], y + Inches(0.03), header_w[2], Inches(0.24),
                     frameworks, font_size=9, color=COLORS["text_light"])
        add_text_box(slide, header_x[3], y + Inches(0.03), header_w[3], Inches(0.24),
                     image, font_size=8, color=COLORS["text_dim"], font_name="Consolas")

    add_text_box(slide, Inches(0.5), Inches(5.15), Inches(9), Inches(0.25),
                 "All images pre-built and pushed to Nexus for air-gapped DinD execution  |  LLM can generate for any language beyond this list",
                 font_size=9, color=COLORS["text_dim"])


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    print("Building slides...")
    build_title_slide(prs)          # 1
    print("  [OK] Title slide")
    build_overview_slide(prs)       # 2
    print("  [OK] Platform overview")
    build_architecture_slide(prs)   # 3
    print("  [OK] Architecture")
    build_frontend_slide(prs)       # 4
    print("  [OK] Frontend dashboard")
    build_demo_flow_slide(prs)      # 5
    print("  [OK] Demo flow")
    build_gitlab_slide(prs)         # 6
    print("  [OK] GitLab deep dive")
    build_jenkins_slide(prs)        # 7
    print("  [OK] Jenkins deep dive")
    build_github_slide(prs)         # 8
    print("  [OK] GitHub Actions deep dive")
    build_rag_learning_slide(prs)   # 9
    print("  [OK] RAG learning loop")
    build_validation_slide(prs)     # 10
    print("  [OK] Validation & self-healing")
    build_languages_slide(prs)      # 11
    print("  [OK] Languages & frameworks")
    build_live_demo_slide(prs)      # 12
    print("  [OK] Live demo")
    build_comparison_slide(prs)     # 13
    print("  [OK] Platform comparison")
    build_closing_slide(prs)        # 14
    print("  [OK] Closing")

    output_path = os.path.join(os.path.dirname(__file__),
                               "AI_DevOps_Pipeline_Demo_v2.pptx")
    prs.save(output_path)
    print(f"\nPresentation saved to: {output_path}")
    print(f"   {len(prs.slides)} slides with animations")


if __name__ == "__main__":
    main()
